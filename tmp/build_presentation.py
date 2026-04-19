"""
Swiftulin — Final Presentation Generator
Mirrors the visual structure of the Road.io reference PPTX.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt
import copy

# ── Color palette ──────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x2A)   # deep background
TEAL      = RGBColor(0x00, 0xB4, 0xD8)   # accent / headings
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF8, 0xFF)   # body background
MUTED     = RGBColor(0xB0, 0xC4, 0xDE)   # secondary text
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

# ── Slide dimensions (widescreen 16:9) ─────────────────────────
W = Inches(10)
H = Inches(5.625)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def add_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, l, t, w, h, fill=None, line=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, text, l, t, w, h,
          size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def bullet_box(slide, items, l, t, w, h, size=16, color=WHITE, heading=None):
    """Text box with optional heading + bullet list."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    if heading:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = heading
        run.font.size = Pt(size + 2)
        run.font.bold = True
        run.font.color.rgb = TEAL

    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.level = 1 if item.startswith("  ") else 0
        run = p.add_run()
        run.text = ("  • " if p.level == 0 else "      ◦ ") + item.strip()
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def accent_line(slide, y):
    """Thin horizontal teal rule."""
    line = slide.shapes.add_shape(1, Inches(0.4), y, Inches(9.2), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def slide_num(slide, n):
    txbox(slide, str(n),
          W - Inches(0.5), H - Inches(0.35),
          Inches(0.4), Inches(0.3),
          size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  BUILD SLIDES
# ══════════════════════════════════════════════════════════════

prs = new_prs()

# ─── Slide 1 ── Title slide ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # BLANK
add_bg(slide, NAVY)

# Full-width teal header bar
bar = add_shape(slide, 0, 0, W, Inches(1.4), fill=TEAL)

# App name in bar
txbox(slide, "Swiftulin",
      Inches(0.5), Inches(0.15),
      Inches(9), Inches(1.0),
      size=52, bold=True, color=WHITE)

# Tagline
txbox(slide, "A Voice-First Insulin Dosing Assistant for Diabetics",
      Inches(0.5), Inches(1.6),
      Inches(9), Inches(0.7),
      size=22, color=TEAL, italic=True)

# Sub-detail
txbox(slide,
      "DS 440 — Capstone Project   |   Penn State University   |   Spring 2026",
      Inches(0.5), Inches(2.4),
      Inches(9), Inches(0.5),
      size=15, color=MUTED)

# Bottom pill
pill = add_shape(slide, Inches(0.5), Inches(4.6), Inches(4.5), Inches(0.65),
                 fill=RGBColor(0x00, 0x4E, 0x64))
txbox(slide, "Gabriel Golden  •  gmg5651@psu.edu",
      Inches(0.6), Inches(4.65), Inches(4.3), Inches(0.55),
      size=13, color=LIGHT_BG)

slide_num(slide, 1)


# ─── Slide 2 ── Problem Statement ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "The Problem",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

bullet_box(slide, [
    "37 million Americans live with diabetes and require daily insulin calculations.",
    "Dosing errors are dangerous — over- or under-dosing can cause seizures or hospitalization.",
    "Existing apps require manual gram entry, which is slow, error-prone, and impractical at meals.",
    "No mainstream solution combines voice logging + real nutritional data + transparent math.",
], Inches(0.5), Inches(1.3), Inches(9), Inches(3.2), size=17)

# Quote callout
box = add_shape(slide, Inches(0.5), Inches(4.3), Inches(8.5), Inches(0.9),
                fill=RGBColor(0x00, 0x4E, 0x64))
txbox(slide,
      '"I just want to say what I ate and know my dose — not look up nutrition labels."',
      Inches(0.65), Inches(4.35), Inches(8.2), Inches(0.8),
      size=14, italic=True, color=LIGHT_BG)

slide_num(slide, 2)


# ─── Slide 3 ── Our Solution ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Our Solution — Swiftulin",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

txbox(slide, "Speak your meal → Get your dose. Instantly. Privately.",
      Inches(0.5), Inches(1.2), Inches(9), Inches(0.55),
      size=18, color=MUTED, italic=True)

# Three feature cards
card_data = [
    ("🎙 Voice Input",       "Say 'I had two eggs and toast' — Swiftulin transcribes and understands."),
    ("🔬 USDA Accuracy",     "Custom NLP pipeline extracts portions; USDA FoodData Central verifies grams."),
    ("💉 Safe Dose Math",    "ICR + ISF formula displayed step-by-step. No hidden calculations, ever."),
]
for i, (title, body) in enumerate(card_data):
    x = Inches(0.3 + i * 3.2)
    card = add_shape(slide, x, Inches(1.9), Inches(3.0), Inches(2.4),
                     fill=RGBColor(0x04, 0x2C, 0x43))
    txbox(slide, title,  x + Inches(0.15), Inches(2.0),  Inches(2.7), Inches(0.5),
          size=14, bold=True, color=TEAL)
    txbox(slide, body,   x + Inches(0.15), Inches(2.55), Inches(2.7), Inches(1.6),
          size=12.5, color=LIGHT_BG)

txbox(slide, "Cross-platform: iOS  •  Android  •  Web (Progressive Web App)",
      Inches(0.5), Inches(4.55), Inches(9), Inches(0.4),
      size=13, color=MUTED, align=PP_ALIGN.CENTER)

slide_num(slide, 3)


# ─── Slide 4 ── Target Users / Problem Scope ──────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Who Is This For?",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

# Left column
txbox(slide, "Primary Users",
      Inches(0.4), Inches(1.2), Inches(4.2), Inches(0.5),
      size=18, bold=True, color=TEAL)
bullet_box(slide, [
    "Adults with Type 1 or insulin-dependent Type 2 diabetes",
    "Users who inject insulin with each meal (MDI regimen)",
    "Caregivers managing doses for children or elderly patients",
], Inches(0.4), Inches(1.7), Inches(4.2), Inches(2.2), size=15)

# Right column
txbox(slide, "Key Pain Points Addressed",
      Inches(5.2), Inches(1.2), Inches(4.4), Inches(0.5),
      size=18, bold=True, color=TEAL)
bullet_box(slide, [
    "Fatigue from manual carb counting at every meal",
    "Fear of dosing mistakes from imprecise estimates",
    "Slow multi-step workflows in existing apps",
    "No privacy guarantee — data shared with third parties",
], Inches(5.2), Inches(1.7), Inches(4.4), Inches(2.6), size=15)

# Divider
div = add_shape(slide, Inches(4.85), Inches(1.2), Pt(2), Inches(3.2), fill=TEAL)

# Stat bar
bar = add_shape(slide, 0, Inches(4.3), W, Inches(0.9), fill=RGBColor(0x00, 0x4E, 0x64))
txbox(slide, "37M Americans with diabetes  •  ~6M on daily insulin  •  100B+ insulin injections per year",
      Inches(0.5), Inches(4.38), Inches(9), Inches(0.6),
      size=13, color=LIGHT_BG, align=PP_ALIGN.CENTER)

slide_num(slide, 4)


# ─── Slide 5 ── Alternative Solutions ─────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Alternative Solutions",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

alts = [
    ("MyFitnessPal",     "Large food database, but no insulin calculation and no voice input."),
    ("Glucose Buddy",    "Tracks readings/logs, but manual entry only. No NLP or USDA lookup."),
    ("OmniPod / Loop",   "Hardware-paired. Requires specific pump. Not universally applicable."),
    ("ChatGPT alone",    "Can estimate carbs, but no USDA verification, no dose math, no privacy."),
]
for i, (name, note) in enumerate(alts):
    y = Inches(1.25 + i * 0.82)
    row = add_shape(slide, Inches(0.4), y, Inches(9.2), Inches(0.7),
                    fill=RGBColor(0x04, 0x2C, 0x43))
    txbox(slide, name, Inches(0.55), y + Inches(0.08), Inches(2.2), Inches(0.55),
          size=15, bold=True, color=TEAL)
    txbox(slide, note, Inches(2.9), y + Inches(0.08), Inches(6.5), Inches(0.55),
          size=14, color=LIGHT_BG)

# Our edge
box = add_shape(slide, Inches(0.4), Inches(4.55), Inches(9.2), Inches(0.7),
                fill=RGBColor(0x00, 0x7A, 0xA0))
txbox(slide, "✅  Swiftulin: Voice + USDA precision + transparent dose math + full privacy — all in one app.",
      Inches(0.55), Inches(4.62), Inches(9.0), Inches(0.55),
      size=14, bold=True, color=WHITE)

slide_num(slide, 5)


# ─── Slide 6 ── How It's Different ────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "How Are We Different?",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

# Comparison table header
cols = ["Feature", "Swiftulin", "Competitors"]
col_w = [Inches(3.5), Inches(2.5), Inches(2.5)]
xs = [Inches(0.5), Inches(4.1), Inches(6.7)]

hdr_y = Inches(1.2)
for i, (col, cw, cx) in enumerate(zip(cols, col_w, xs)):
    hdr = add_shape(slide, cx, hdr_y, cw - Inches(0.05), Inches(0.45),
                    fill=TEAL)
    txbox(slide, col, cx + Inches(0.05), hdr_y + Inches(0.05),
          cw - Inches(0.1), Inches(0.4),
          size=13, bold=True, color=NAVY)

rows = [
    ("Voice-to-carb pipeline",   "✅ Full",     "❌ Manual entry"),
    ("NLP + USDA verification",  "✅ Regex + API", "❌ Estimate only"),
    ("Insulin dose calculator",  "✅ With math", "⚠️ Partial / none"),
    ("Privacy — no data logs",   "✅ Ephemeral", "❌ Data uploaded"),
    ("Works offline (core)",     "✅ SQLite",    "⚠️ Partial"),
    ("Cross-platform",           "✅ iOS/Android/Web", "⚠️ Varies"),
]
row_colors = [RGBColor(0x04, 0x2C, 0x43), RGBColor(0x05, 0x35, 0x50)]
for r, (feat, us, them) in enumerate(rows):
    ry = Inches(1.7 + r * 0.48)
    rc = row_colors[r % 2]
    for i, (text, cx, cw) in enumerate(zip([feat, us, them], xs, col_w)):
        cell = add_shape(slide, cx, ry, cw - Inches(0.05), Inches(0.44), fill=rc)
        txbox(slide, text, cx + Inches(0.05), ry + Inches(0.05),
              cw - Inches(0.1), Inches(0.4),
              size=12.5, color=LIGHT_BG if i == 0 else (TEAL if "✅" in text else MUTED))

slide_num(slide, 6)


# ─── Slide 7 ── Technical Architecture ────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Technical Architecture",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

# Pipeline boxes
pipeline = [
    ("🎙", "Voice\nCapture",     "Expo Voice\nModule"),
    ("🤖", "NLP\nExtraction",   "Custom Regex\nParser"),
    ("🥗", "Nutritional\nLookup","USDA Food-\nData Central"),
    ("💊", "Dose\nCalculation", "ICR + ISF\nFormula"),
    ("✅", "User\nConfirmation", "Manual tap\nrequired"),
]
box_w = Inches(1.7)
gap   = Inches(0.2)
start_x = Inches(0.35)

for i, (icon, title, sub) in enumerate(pipeline):
    bx = start_x + i * (box_w + gap)
    pbox = add_shape(slide, bx, Inches(1.35), box_w, Inches(1.85),
                     fill=RGBColor(0x04, 0x2C, 0x43))
    txbox(slide, icon,  bx + Inches(0.1), Inches(1.4),  box_w - Inches(0.2), Inches(0.4),
          size=22, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(slide, title, bx + Inches(0.05), Inches(1.82), box_w - Inches(0.1), Inches(0.5),
          size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txbox(slide, sub,   bx + Inches(0.05), Inches(2.32), box_w - Inches(0.1), Inches(0.75),
          size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < len(pipeline) - 1:
        ax = bx + box_w + Inches(0.02)
        txbox(slide, "→", ax, Inches(1.9), gap + Inches(0.1), Inches(0.4),
              size=16, color=TEAL, align=PP_ALIGN.CENTER)

# Stack chips
txbox(slide, "Tech Stack:",
      Inches(0.5), Inches(3.4), Inches(2), Inches(0.4),
      size=14, bold=True, color=MUTED)

chips = [
    "React Native / Expo", "Expo Router", "TypeScript",
    "Zustand", "React Query", "Drizzle ORM + SQLite",
    "EAS Build (iOS/Android)", "Jest + Maestro E2E",
]
chip_y = Inches(3.85)
cx = Inches(0.4)
for chip in chips:
    cw = Inches(len(chip) * 0.098 + 0.3)
    if cx + cw > Inches(9.8):
        cx = Inches(0.4)
        chip_y += Inches(0.45)
    c = add_shape(slide, cx, chip_y, cw, Inches(0.38),
                  fill=RGBColor(0x00, 0x4E, 0x64))
    txbox(slide, chip, cx + Inches(0.07), chip_y + Inches(0.04),
          cw - Inches(0.1), Inches(0.3),
          size=10.5, color=LIGHT_BG)
    cx += cw + Inches(0.12)

slide_num(slide, 7)


# ─── Slide 8 ── Features Overview ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Core Features",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

feats = [
    ("F-001", "Voice-to-Text Input",
     "Tap to record • Waveform visualizer • Editable transcription fallback"),
    ("F-002", "NLP + USDA Lookup",
     "Regex extracts food/quantity • USDA verifies grams • USDA/Estimate badges"),
    ("F-003", "Insulin Calculation",
     "Meal dose (Carbs ÷ ICR) + Correction dose • Full math shown • Safety caps"),
    ("F-004", "Manual Entry",
     "Type instead of speaking • Supports decimals and fractions (0.1 step)"),
    ("F-005", "Medical Settings",
     "ICR, ISF, target BG stored in encrypted SQLite • No cloud upload"),
    ("F-006", "Offline Infrastructure",
     ".native.ts / .ts platform split • SQLite on device • localStorage on web"),
]
for i, (fid, name, detail) in enumerate(feats):
    r, c = divmod(i, 2)
    fx = Inches(0.4 + c * 4.75)
    fy = Inches(1.25 + r * 1.25)
    fbox = add_shape(slide, fx, fy, Inches(4.55), Inches(1.1),
                     fill=RGBColor(0x04, 0x2C, 0x43))
    txbox(slide, f"{fid}  {name}", fx + Inches(0.12), fy + Inches(0.08),
          Inches(4.3), Inches(0.38),
          size=13, bold=True, color=TEAL)
    txbox(slide, detail, fx + Inches(0.12), fy + Inches(0.48),
          Inches(4.3), Inches(0.55),
          size=11.5, color=MUTED)

slide_num(slide, 8)


# ─── Slide 9 ── Demo — Voice Flow ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Demonstration — Voice Logging Flow",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

txbox(slide, 'User says: "I had a cup of oatmeal with a banana and a glass of OJ"',
      Inches(0.5), Inches(1.25), Inches(9), Inches(0.55),
      size=16, italic=True, color=MUTED)

steps = [
    ("1", "Tap 🎙",                "HomeStartScreen → recording starts, waveform animates"),
    ("2", "Transcription",          "Expo Voice → text appears in editable field"),
    ("3", "NLP Extraction",         "Parser extracts: oatmeal 1 cup, banana 1 med, OJ 8 oz"),
    ("4", "USDA Verification",      "Each food fetched from USDA FoodData Central; carbs confirmed"),
    ("5", "Results Screen",         "List with USDA badges; ± 0.1 adjusters; total carbs displayed"),
    ("6", "Dose Calculation",       "User taps 'Calculate Dose'; modal shows ICR + correction math"),
    ("7", "Confirm & Done",         "User taps 'Complete'; app resets (no data saved)"),
]
col_w_s = [Inches(0.45), Inches(2.0), Inches(6.2)]
xs_s    = [Inches(0.4),  Inches(0.9), Inches(3.0)]
for r, (num, steptitle, detail) in enumerate(steps):
    ry = Inches(1.95 + r * 0.46)
    row_fill = RGBColor(0x04, 0x2C, 0x43) if r % 2 == 0 else RGBColor(0x05, 0x35, 0x50)
    add_shape(slide, Inches(0.4), ry, Inches(9.2), Inches(0.43), fill=row_fill)
    txbox(slide, num,       Inches(0.5),  ry+Inches(0.04), Inches(0.38), Inches(0.38),
          size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txbox(slide, steptitle, Inches(0.95), ry+Inches(0.04), Inches(1.9),  Inches(0.38),
          size=12.5, bold=True, color=WHITE)
    txbox(slide, detail,    Inches(2.95), ry+Inches(0.04), Inches(6.5),  Inches(0.38),
          size=12, color=MUTED)

slide_num(slide, 9)


# ─── Slide 10 ── Demo — Manual Entry & Settings ───────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Demonstration — Manual Entry & Settings",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

# Left: Manual entry
txbox(slide, "Manual Entry Flow",
      Inches(0.4), Inches(1.2), Inches(4.3), Inches(0.45),
      size=18, bold=True, color=TEAL)
bullet_box(slide, [
    'Tap "Manual Entry" on home screen',
    "Type food description or search by name",
    "Adjust gram weight with ± 0.1 stepper",
    "Add multiple items then calculate dose",
], Inches(0.4), Inches(1.7), Inches(4.3), Inches(2.0), size=14)

# Divider
add_shape(slide, Inches(4.84), Inches(1.2), Pt(2), Inches(3.2), fill=TEAL)

# Right: Settings
txbox(slide, "Medical Settings Profile",
      Inches(5.1), Inches(1.2), Inches(4.5), Inches(0.45),
      size=18, bold=True, color=TEAL)
bullet_box(slide, [
    "Insulin-to-Carb Ratio (ICR) — e.g., 1 unit per 10g",
    "Insulin Sensitivity Factor (ISF) — e.g., 1 unit drops BG 40 mg/dL",
    "Target blood glucose — e.g., 100 mg/dL",
    "Stored in encrypted SQLite — never leaves device",
], Inches(5.1), Inches(1.7), Inches(4.5), Inches(2.3), size=14)

# Bottom formula box
fbox = add_shape(slide, Inches(0.4), Inches(4.35), Inches(9.2), Inches(0.85),
                 fill=RGBColor(0x00, 0x4E, 0x64))
txbox(slide,
      "Dose = (Total Carbs ÷ ICR)  +  ((Current BG − Target BG) ÷ ISF)",
      Inches(0.55), Inches(4.42), Inches(9.0), Inches(0.7),
      size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_num(slide, 10)


# ─── Slide 11 ── Privacy & Safety ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Privacy & Safety Design",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

pillars = [
    ("🔒  Ephemeral by Design",
     "No meal logs, no insulin history. Every session is stateless. What you eat stays with you."),
    ("🌍  Offline-Capable NLP",
     "Regex parsing runs entirely on-device without requiring LLM calls or internet connection."),
    ("📐  Transparent Math",
     "Every insulin calculation shows the exact formula used. No black-box dose recommendations."),
    ("✋  Human in the Loop",
     "Every automated result requires a manual 'Complete' tap before any action is taken."),
]
for i, (title, body) in enumerate(pillars):
    r, c = divmod(i, 2)
    px = Inches(0.4 + c * 4.75)
    py = Inches(1.25 + r * 1.55)
    pbox = add_shape(slide, px, py, Inches(4.55), Inches(1.35),
                     fill=RGBColor(0x04, 0x2C, 0x43))
    txbox(slide, title, px + Inches(0.12), py + Inches(0.1),
          Inches(4.3), Inches(0.45),
          size=14, bold=True, color=TEAL)
    txbox(slide, body,  px + Inches(0.12), py + Inches(0.55),
          Inches(4.3), Inches(0.7),
          size=12.5, color=LIGHT_BG)

slide_num(slide, 11)


# ─── Slide 12 ── Testing & Quality ────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Testing & Quality",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

# Left: test suite
txbox(slide, "21 Passing Test Cases",
      Inches(0.4), Inches(1.2), Inches(4.3), Inches(0.45),
      size=18, bold=True, color=TEAL)

tests = [
    ("insulinCalculator.test.ts",  "Dose math for ICR + ISF + corrections"),
    ("nlpExtractor.test.ts",       "Food and quantity parse accuracy"),
    ("useVoiceToText.test.ts",     "Recording hook state transitions"),
    ("settingsStore.test.ts",      "ICR/ISF persistence and retrieval"),
    ("ManualCarbModal.test.tsx",   "Manual entry form validation"),
    ("HomeRecordingScreen.test",   "UI render + recording interactions"),
    ("ProfileSettingsScreen.test", "Settings save/load round-trip"),
    ("db.test.ts",                 "SQLite schema + migrations"),
]
for i, (file, desc) in enumerate(tests):
    ty = Inches(1.72 + i * 0.38)
    add_shape(slide, Inches(0.4), ty, Inches(4.3), Inches(0.36),
              fill=RGBColor(0x04, 0x2C, 0x43) if i % 2 == 0 else RGBColor(0x05, 0x35, 0x50))
    txbox(slide, file, Inches(0.5), ty + Inches(0.03), Inches(2.3), Inches(0.32),
          size=10.5, bold=True, color=TEAL)
    txbox(slide, desc, Inches(2.85), ty + Inches(0.03), Inches(1.8), Inches(0.32),
          size=10.5, color=MUTED)

# Divider
add_shape(slide, Inches(4.84), Inches(1.2), Pt(2), Inches(3.7), fill=TEAL)

# Right: methodology
txbox(slide, "Quality Methodology",
      Inches(5.1), Inches(1.2), Inches(4.5), Inches(0.45),
      size=18, bold=True, color=TEAL)
bullet_box(slide, [
    "Test-Driven Development (TDD) — tests written before implementation",
    "Jest for unit and integration testing",
    "Maestro E2E YAML spec for full user flow",
    "TypeScript strict mode across all files",
    "EAS Build CI/CD — automated cloud builds on merge",
    "Coverage reports generated on every test run",
], Inches(5.1), Inches(1.7), Inches(4.5), Inches(2.9), size=13.5)

slide_num(slide, 12)


# ─── Slide 13 ── Why Is This Valuable? ────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Why Is This Valuable?",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

vals = [
    ("For the Patient",
     "Faster, safer dosing with fewer mental steps. Reduces cognitive load at every meal."),
    ("For Safety",
     "Math transparency prevents blind trust in a number. Every dose is auditable by the user."),
    ("For Privacy",
     "Sensitive health data never leaves the device. No ads, no profiles, no sharing."),
    ("For Accessibility",
     "Voice-first design helps users with dexterity issues or visual impairments."),
    ("For Scale",
     "Expo's universal build targets iOS, Android, and Web from one codebase."),
    ("For Research",
     "Structured logging architecture ready to plug in anonymized datasets for future ML."),
]
for i, (title, body) in enumerate(vals):
    r, c = divmod(i, 2)
    vx = Inches(0.4 + c * 4.75)
    vy = Inches(1.25 + r * 1.2)
    vbox = add_shape(slide, vx, vy, Inches(4.55), Inches(1.05),
                     fill=RGBColor(0x04, 0x2C, 0x43))
    left_bar = add_shape(slide, vx, vy, Inches(0.08), Inches(1.05), fill=TEAL)
    txbox(slide, title, vx + Inches(0.18), vy + Inches(0.07),
          Inches(4.25), Inches(0.35),
          size=13, bold=True, color=TEAL)
    txbox(slide, body,  vx + Inches(0.18), vy + Inches(0.45),
          Inches(4.25), Inches(0.52),
          size=12, color=LIGHT_BG)

slide_num(slide, 13)


# ─── Slide 14 ── Architectural Decisions ──────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Key Architectural Decisions",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

decisions = [
    ("Drizzle ORM + SQLite",         "Type-safe schema, migrations, and query builder over raw SQLite."),
    ("Zustand over Redux",            "Minimal boilerplate, hooks-first API — no Provider needed."),
    ("Expo Router",                  "File-based routing with first-class Expo support and web compat."),
    ("React Query for remote data",  "Built-in async state, caching, and offline queue."),
    ("EAS Internal Distribution",    "Secure, repeatable cloud builds for iOS/Android preview."),
    ("Platform-Isolated Architecture","`.native.ts` extensions strictly separate web mocks from native modules."),
    ("Regex + USDA dual-layer NLP",  "Custom NLP engine for extraction; USDA for ground-truth carb data."),
]
for i, (decision, rationale) in enumerate(decisions):
    dy = Inches(1.22 + i * 0.48)
    fill = RGBColor(0x04, 0x2C, 0x43) if i % 2 == 0 else RGBColor(0x05, 0x35, 0x50)
    add_shape(slide, Inches(0.4), dy, Inches(9.2), Inches(0.44), fill=fill)
    txbox(slide, decision,  Inches(0.55), dy + Inches(0.04), Inches(3.3), Inches(0.38),
          size=12.5, bold=True, color=TEAL)
    txbox(slide, rationale, Inches(3.95), dy + Inches(0.04), Inches(5.5), Inches(0.38),
          size=12.5, color=MUTED)

slide_num(slide, 14)


# ─── Slide 15 ── Challenges & Solutions ───────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Challenges & How We Solved Them",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

challenges = [
    ("NLP Ambiguity in Food Descriptions",
     "Fraction support (e.g. 3/4 cup), unit normalization, multi-food extraction in one utterance.",
     "Regex segmentation layer with rigorous fractional parsing + USDA fallback ranks."),
    ("Cross-Platform SQLite",
     "SQLite not available in web browsers; needed unified database API.",
     "Platform-isolated `.native.ts` / `.ts` extensions with shared Drizzle schema."),
    ("Edge Computing",
     "Need for instant, offline-capable food extraction without API latency.",
     "Swapped LLM design for a high-performance, on-device regex NLP pipeline."),
    ("Dosing Safety",
     "Any calculation error in a medical app can cause real harm.",
     "Hard max values, transparent math display, and mandatory manual confirmation before dose."),
]
for i, (chall, prob, sol) in enumerate(challenges):
    cy = Inches(1.2 + i * 0.96)
    add_shape(slide, Inches(0.4), cy, Inches(9.2), Inches(0.9),
              fill=RGBColor(0x04, 0x2C, 0x43) if i % 2 == 0 else RGBColor(0x05, 0x35, 0x50))
    txbox(slide, f"⚠ {chall}", Inches(0.55), cy + Inches(0.04), Inches(3.5), Inches(0.38),
          size=12.5, bold=True, color=TEAL)
    txbox(slide, f"Problem: {prob}", Inches(0.55), cy + Inches(0.44), Inches(4.1), Inches(0.38),
          size=11, color=MUTED, italic=True)
    txbox(slide, f"✅ {sol}", Inches(4.75), cy + Inches(0.04), Inches(4.7), Inches(0.8),
          size=11.5, color=LIGHT_BG)

slide_num(slide, 15)


# ─── Slide 16 ── Conclusion ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

txbox(slide, "Conclusion",
      Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
      size=34, bold=True, color=TEAL)
accent_line(slide, Inches(1.1))

txbox(slide,
      "Swiftulin delivers a complete, privacy-first insulin dosing pipeline that no existing app provides.",
      Inches(0.5), Inches(1.2), Inches(9), Inches(0.6),
      size=16, color=LIGHT_BG, italic=True)

summary = [
    "✅  Voice-to-carb pipeline powered by custom NLP + USDA FoodData Central",
    "✅  Transparent insulin dose math — ICR + ISF shown step by step",
    "✅  Zero data persistence — fully ephemeral, privacy-first architecture",
    "✅  21 passing tests with TDD methodology and TypeScript strict mode",
    "✅  Cross-platform (iOS, Android, Web) from a single Expo codebase",
    "✅  Production-grade CI/CD with EAS and platform-isolated build system",
]
for i, line in enumerate(summary):
    ly = Inches(1.95 + i * 0.42)
    txbox(slide, line, Inches(0.7), ly, Inches(8.5), Inches(0.38),
          size=15, color=LIGHT_BG if "✅" not in line else WHITE)

txbox(slide,
      "Future Work: CGM integration • Meal history analytics (opt-in) • Smartwatch companion",
      Inches(0.5), Inches(4.55), Inches(9), Inches(0.45),
      size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

slide_num(slide, 16)


# ─── Slide 17 ── Thank You ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

# Teal bar
add_shape(slide, 0, Inches(1.8), W, Inches(2.1), fill=RGBColor(0x00, 0x4E, 0x64))

txbox(slide, "Thank You!",
      Inches(0.5), Inches(2.0), Inches(9), Inches(0.9),
      size=48, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

txbox(slide, "Questions? Feedback?",
      Inches(0.5), Inches(2.95), Inches(9), Inches(0.5),
      size=18, color=WHITE, align=PP_ALIGN.CENTER)

txbox(slide, "Swiftulin",
      Inches(0.5), Inches(0.4), Inches(9), Inches(0.6),
      size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txbox(slide, "DS 440 Capstone  •  Penn State University  •  Spring 2026",
      Inches(0.5), Inches(0.95), Inches(9), Inches(0.4),
      size=13, color=MUTED, align=PP_ALIGN.CENTER)

txbox(slide, "Gabriel Golden  •  gmg5651@psu.edu",
      Inches(0.5), Inches(4.4), Inches(9), Inches(0.5),
      size=14, color=MUTED, align=PP_ALIGN.CENTER)

txbox(slide, "github.com/gmg5651/DS440Project",
      Inches(0.5), Inches(4.9), Inches(9), Inches(0.4),
      size=13, color=TEAL, align=PP_ALIGN.CENTER)

slide_num(slide, 17)


# ── Save ───────────────────────────────────────────────────────
out_path = "/Users/gabegolden/Desktop/Swiftulin_Final_Presentation.pptx"
prs.save(out_path)
print(f"✅  Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
